"""
Complete Self-Healing RAG System — FAISS + PyPDF edition

Integrates all components: HyDE, Query Decomposition, CRAG, Cross-Encoder Reranking,
and Dynamic Learning into a unified self-healing RAG system.

Vector store: FAISS (via langchain_community)
Embeddings:   HuggingFaceEmbeddings  (BAAI/bge-small-en-v1.5)
Chunking:     RecursiveCharacterTextSplitter
PDF parsing:  pypdf
"""

import os
import tempfile
from typing import List, Dict, Any, Optional
from datetime import datetime
from types import SimpleNamespace

from dotenv import load_dotenv

# LangChain core
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_nvidia_ai_endpoints import ChatNVIDIA

# PDF parsing
from pypdf import PdfReader

# LLM bridge (LlamaIndex LLM interface kept for NVIDIA calls)
from llama_index.llms.langchain import LangChainLLM
from llama_index.embeddings.langchain import LangchainEmbedding
from llama_index.core import Settings

# Self-healing components
from hyde import build_hyde_retriever
from query_decomposition import QueryDecomposer
from corrective_rag import CRAGSystem
from reranker import Reranker
from dynamic_prompting import LearningManager, PromptOptimizationTracker
from intent_analyzer import IntentAnalyzer


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
EMBED_MODEL_NAME = "BAAI/bge-small-en-v1.5"
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50


class SelfHealingRAGSystem:
    """
    Complete self-healing RAG system with all advanced techniques.

    Document ingestion pipeline:
        PDF → pypdf → RecursiveCharacterTextSplitter
            → HuggingFaceEmbeddings → FAISS

    Query pipeline:
        Query → (Decomposition) → (HyDE) → FAISS retrieval
             → CRAG validation → Cross-Encoder reranking
             → Dynamic few-shot learning → LLM answer generation
    """

    def __init__(
        self,
        nvidia_api_key: str = None,
        tavily_api_key: str = None,
        enable_web_search: bool = False,
    ):
        """
        Initialize RAG system.

        Args:
            nvidia_api_key: NVIDIA API key (also read from env).
            tavily_api_key: Tavily API key for web-search fallback.
            enable_web_search: Whether to enable web-search in CRAG.
        """
        load_dotenv(override=True)
        print("🚀 Initializing Self-Healing RAG System...")

        # ── LLM setup ────────────────────────────────────────────────────────
        try:
            nvidia_model = "meta/llama-3.1-8b-instruct"
            chat_client = ChatNVIDIA(
                model=nvidia_model,
                temperature=0.3,
                api_key=os.getenv("NVIDIA_API_KEY"),
            )
            self.llm = LangChainLLM(llm=chat_client)
            Settings.llm = self.llm
        except Exception as e:
            print(f"⚠️  NVIDIA LLM init failed: {e}. Using local stub LLM.")

            class _FallbackLLM:
                """Minimal stub so the pipeline can run without a real LLM."""

                def complete(self, prompt: str):
                    text = (
                        "(LLM fallback) NVIDIA LLM unavailable. "
                        "Please check your NVIDIA_API_KEY.\n\n"
                        "Prompt received (truncated): " + prompt[:200] + "..."
                    )
                    return SimpleNamespace(text=text)

            self.llm = _FallbackLLM()
            Settings.llm = self.llm

        # ── Embedding model ───────────────────────────────────────────────────
        print("  📦 Loading embedding model…")
        self.embeddings = HuggingFaceEmbeddings(model_name=EMBED_MODEL_NAME)
        Settings.embed_model = LangchainEmbedding(self.embeddings)

        # ── Text splitter ─────────────────────────────────────────────────────
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

        # ── Self-healing components ───────────────────────────────────────────
        print("  📦 Loading components...")
        self.query_decomposer = QueryDecomposer(model="meta/llama-3.1-8b-instruct")
        self.reranker = Reranker()
        self.learning_manager = LearningManager(similarity_top_k=2)
        self.crag_system = CRAGSystem(
            grader_model="meta/llama-3.1-8b-instruct",
            generator_model="meta/llama-3.1-8b-instruct",
            web_search_enabled=enable_web_search,
            tavily_api_key=tavily_api_key,
        )
        self.optimization_tracker = PromptOptimizationTracker()
        self.intent_analyzer = IntentAnalyzer(model="meta/llama-3.1-8b-instruct")

        # ── State ─────────────────────────────────────────────────────────────
        self.vectorstore: Optional[FAISS] = None
        self.hyde_retriever = None
        self.loaded_chunks: int = 0



        self.session_history: Dict[str, List[Dict[str, str]]] = {}

        print("✅ System initialized successfully!")

    # ── Document ingestion ────────────────────────────────────────────────────

    def _build_vectorstore(self, chunks: List[Document]) -> None:
        """Build (or extend) FAISS vectorstore from chunks."""
        if self.vectorstore is None:
            self.vectorstore = FAISS.from_documents(chunks, self.embeddings)
        else:
            self.vectorstore.add_documents(chunks)

        self.loaded_chunks += len(chunks)

        # Rebuild HyDE retriever with updated store
        self.hyde_retriever = build_hyde_retriever(
            vectorstore=self.vectorstore,
            llm=self.llm,
            top_k=5,
            include_original=True,
        )



    def load_file(self, file_path: str, original_filename: str = None) -> int:
        """
        Smart parser to load PDF, DOCX, PPTX, HTML, or TXT files.

        Args:
            file_path: Absolute path to the file.
            original_filename: Original filename (useful if file_path is a tmp file).

        Returns:
            Number of chunks indexed.
        """
        filename = original_filename or os.path.basename(file_path)
        ext = filename.lower().split('.')[-1]
        print(f"📄 Parsing {ext.upper()} file: {filename}")

        raw_docs: List[Document] = []

        try:
            if ext == "pdf":
                from langchain_community.document_loaders import PyPDFLoader
                loader = PyPDFLoader(file_path)
                docs = loader.load()
                # Update source metadata to original filename
                for doc in docs:
                    doc.metadata['source'] = filename
                raw_docs.extend(docs)
            elif ext in ["html", "htm"]:
                from langchain_community.document_loaders import BSHTMLLoader
                loader = BSHTMLLoader(file_path, open_encoding="utf-8")
                docs = loader.load()
                for doc in docs:
                    doc.metadata['source'] = filename
                raw_docs.extend(docs)
            elif ext == "txt":
                from langchain_community.document_loaders import TextLoader
                loader = TextLoader(file_path, encoding="utf-8")
                docs = loader.load()
                for doc in docs:
                    doc.metadata['source'] = filename
                raw_docs.extend(docs)
            elif ext == "docx":
                import docx
                doc = docx.Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                if text:
                    raw_docs.append(Document(page_content=text, metadata={"source": filename}))
            elif ext == "pptx":
                from pptx import Presentation
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = "\n".join(text_runs).strip()
                if text:
                    raw_docs.append(Document(page_content=text, metadata={"source": filename}))
            else:
                print(f"  ⚠️  Unsupported file format: {ext}")
                return 0
        except Exception as e:
            print(f"  ⚠️  Error parsing {filename}: {e}")
            return 0

        if not raw_docs:
            print(f"  ⚠️  No text found in {filename}.")
            return 0

        print(f"  📖 Extracted {len(raw_docs)} document sections/pages")

        chunks = self.splitter.split_documents(raw_docs)
        self._build_vectorstore(chunks)

        print(f"  ✅ Indexed {len(chunks)} chunks from {filename}")
        return len(chunks)



    # ── Retrieval ─────────────────────────────────────────────────────────────

    @property
    def vector_index(self):
        """Compatibility alias used by api_server.py health checks."""
        return self.vectorstore

    def enhanced_retrieve(
        self,
        query: str,
        use_hyde: bool = True,
        top_k: int = 5,
    ) -> List[Document]:
        """
        Enhanced retrieval with optional HyDE.

        Args:
            query: User query string.
            use_hyde: Whether to use HyDE transformation.
            top_k: Number of documents to return from plain search.

        Returns:
            List of LangChain Document objects.
        """
        if self.vectorstore is None:
            print("  ⚠️  No vectorstore loaded yet.")
            return []

        print(f"🔍 Retrieving for: '{query[:60]}…'")

        if use_hyde and self.hyde_retriever is not None:
            print("  🧠 Using HyDE transformation…")
            docs = self.hyde_retriever.retrieve(query)
        else:
            print("  📖 Using standard FAISS retrieval…")
            docs = self.vectorstore.similarity_search(query, k=top_k)

        print(f"  ✅ Retrieved {len(docs)} documents")
        return docs

    # ── Main query pipeline ───────────────────────────────────────────────────

    def process_query(
        self,
        query: str,
        enable_decomposition: bool = True,
        enable_hyde: bool = True,
        enable_crag: bool = True,
        enable_reranking: bool = True,
        enable_learning: bool = True,
        manual_override: bool = False,
        thread_id: str = "default",
    ) -> Dict[str, Any]:
        """
        Process query through the complete self-healing pipeline.

        Steps:
            1. Query Enhancement (Decomposition + HyDE)
            2. CRAG Document Validation
            3. Cross-Encoder Reranking
            4. Dynamic Few-Shot Learning
            5. Answer Generation

        Returns:
            Dict with answer and rich metadata.
        """
        start_time = datetime.now()

        print(f"\n{'=' * 70}")
        print("🔄 SELF-HEALING RAG PIPELINE")
        print(f"{'=' * 70}")
        print(f"Query: {query}")
        print(f"{'=' * 70}\n")

        techniques_used: List[str] = []

        # ── Step 0: Intent Analysis ───────────────────────────────────────────
        intent = self.intent_analyzer.analyze(query)
        print(f"🎯 INTENT DETECTED: {intent}")

        # If not manual override, the system is fully autonomous and overrides input flags
        if not manual_override:
            enable_decomposition = True
            enable_hyde = True
            enable_crag = True
            enable_reranking = True
            enable_learning = True

        # If it's a global review and we are in autonomous mode, disable HyDE and CRAG
        if intent == "global_review" and not manual_override:
            print("  → Intent is global_review: Disabling HyDE & CRAG, increasing top_k")
            use_hyde = False
            use_crag = False
            retrieval_k = 20
        else:
            use_hyde = enable_hyde
            use_crag = enable_crag
            retrieval_k = 5

        # ── Step 1: Query Enhancement ─────────────────────────────────────────
        sub_queries = [query]
        all_documents: List[Document] = []

        should_decompose = enable_decomposition and (manual_override or (intent == "qa" and self.query_decomposer.is_complex_query(query)))

        if should_decompose:
            print("📊 STEP 1: Query Decomposition")
            print("-" * 70)
            try:
                sub_queries = self.query_decomposer.decompose(query)
                if len(sub_queries) > 1:
                    print(f"  ✓ Decomposed into {len(sub_queries)} sub-queries:")
                    for i, sq in enumerate(sub_queries, 1):
                        print(f"    {i}. {sq}")

                    for sq in sub_queries:
                        docs = self.enhanced_retrieve(sq, use_hyde=use_hyde, top_k=max(3, retrieval_k // len(sub_queries)))
                        all_documents.extend(docs)

                    techniques_used.append("Query Decomposition")
                else:
                    all_documents = self.enhanced_retrieve(query, use_hyde=use_hyde, top_k=retrieval_k)
            except Exception as e:
                print(f"  ⚠ Decomposition error: {e}")
                all_documents = self.enhanced_retrieve(query, use_hyde=use_hyde, top_k=retrieval_k)
        else:
            print("📊 STEP 1: Standard Retrieval")
            print("-" * 70)
            all_documents = self.enhanced_retrieve(query, use_hyde=use_hyde, top_k=retrieval_k)

        if use_hyde:
            techniques_used.append("HyDE")

        print(f"\n  📄 Total documents retrieved: {len(all_documents)}\n")

        # ── Step 2: CRAG Document Validation ─────────────────────────────────
        filtered_docs = all_documents
        if use_crag and all_documents:
            print("🔍 STEP 2: CRAG Document Validation")
            print("-" * 70)
            try:
                result = self.crag_system.run(query, all_documents, thread_id=thread_id)
                filtered_docs = result.get("documents", all_documents)

                removed = len(all_documents) - len(filtered_docs)
                if removed > 0:
                    print(f"  🚨 Filtered out {removed} irrelevant document(s)")
                    techniques_used.append("CRAG")
                else:
                    print("  ✓ All documents passed relevance check")
            except Exception as e:
                print(f"  ⚠ CRAG error: {e}")
                filtered_docs = all_documents

        print(f"\n  📄 Documents after CRAG: {len(filtered_docs)}\n")

        # ── Step 3: Cross-Encoder Reranking ───────────────────────────────────
        final_docs = filtered_docs
        
        should_rerank = enable_reranking and len(filtered_docs) > 1 and (manual_override or intent == "qa")
        if should_rerank:
            print("🎯 STEP 3: Cross-Encoder Reranking")
            print("-" * 70)
            try:
                doc_texts = [doc.page_content for doc in filtered_docs]
                reranked_texts = self.reranker.rerank(query, doc_texts, top_k=min(5, len(doc_texts)))

                final_docs = []
                for text in reranked_texts:
                    for doc in filtered_docs:
                        if doc.page_content == text:
                            final_docs.append(doc)
                            break

                print(f"  ✓ Reranked to top {len(final_docs)} documents")
                techniques_used.append("Cross-Encoder Reranking")
            except Exception as e:
                print(f"  ⚠ Reranking error: {e}")
                final_docs = filtered_docs[:5]
        elif enable_reranking and len(filtered_docs) > 1 and intent == "global_review":
            print("🎯 STEP 3: Skipping Reranking for Global Review (preserving broad context)")
            print("-" * 70)

        print(f"\n  📄 Final documents: {len(final_docs)}\n")

        # ── Step 4: Dynamic Few-Shot Learning ────────────────────────────────
        few_shot_context = ""
        if enable_learning:
            print("🧠 STEP 4: Dynamic Few-Shot Learning")
            print("-" * 70)
            try:
                few_shot_context = self.learning_manager.get_dynamic_prompt(query)
                if few_shot_context:
                    techniques_used.append("Dynamic Learning")
                    print("  ✓ Applied learned examples")
                else:
                    print("  ℹ No relevant past examples found")
            except Exception as e:
                print(f"  ⚠ Learning error: {e}")

        print()

        # ── Step 5: Answer Generation ─────────────────────────────────────────
        print("✍️  STEP 5: Answer Generation")
        print("-" * 70)

        if thread_id not in self.session_history:
            self.session_history[thread_id] = []
        chat_history = self.session_history[thread_id]

        answer = self.generate_answer(query, final_docs, few_shot_context, chat_history)
        print("  ✓ Answer generated\n")

        self.session_history[thread_id].append({"role": "user", "content": query})
        self.session_history[thread_id].append({"role": "assistant", "content": answer})

        # Keep last 5 turns (10 messages)
        if len(self.session_history[thread_id]) > 10:
            self.session_history[thread_id] = self.session_history[thread_id][-10:]

        end_time = datetime.now()
        processing_time = (end_time - start_time).total_seconds()

        result = {
            "query": query,
            "sub_queries": sub_queries,
            "answer": answer,
            "documents_retrieved": len(all_documents),
            "documents_after_crag": len(filtered_docs),
            "final_documents": len(final_docs),
            "document_contents": [
                doc.page_content[:200] + "…"
                for doc in final_docs[:3]
            ],
            "processing_time": round(processing_time, 2),
            "techniques_used": techniques_used,
            "timestamp": datetime.now().isoformat(),
        }

        print(f"{'=' * 70}")
        print(f"✅ Pipeline completed in {processing_time:.2f}s")
        print(f"📊 Techniques: {', '.join(techniques_used)}")
        print(f"{'=' * 70}\n")

        return result

    # ── Answer generation ─────────────────────────────────────────────────────

    def generate_answer(
        self,
        query: str,
        documents: List[Document],
        few_shot_context: str = "",
        chat_history: List[Dict[str, str]] = None,
    ) -> str:
        """
        Generate an answer from retrieved documents.

        Args:
            query: User query.
            documents: Retrieved and validated documents.
            few_shot_context: Optional few-shot examples from dynamic learning.
            chat_history: Conversation history for this session.

        Returns:
            Generated answer string.
        """
        if not documents:
            return "I apologize, but I couldn't find relevant information to answer your question."

        # Combine document content (more docs if global_review)
        max_docs = 20 if len(documents) > 5 else 5
        context = "\n\n".join(
            f"[Document {i+1}]\n{doc.page_content}"
            for i, doc in enumerate(documents[:max_docs])
        )

        prompt_parts: List[str] = []

        if few_shot_context:
            prompt_parts.append(few_shot_context)
            prompt_parts.append("\n---\n")

        prompt_parts.extend([
            "You are a helpful, conversational technical assistant.",
            "Use the conversation history for context but base factual answers on the provided documents.",
            "If the documents don't contain enough information, say so clearly.\n",
            "CRITICAL INSTRUCTIONS:",
            "- Keep titles bold and format them neatly using Markdown.",
            "- Add meaningful emojis throughout your response to make it engaging.",
            "- If making a comparison, you MUST create a neat Markdown table.",
            "- Be very detailed, broad, and cover every aspect of the question.",
            "- Provide a long, comprehensive, and exhaustive answer.\n",
        ])

        if chat_history:
            prompt_parts.append("Conversation History:")
            for msg in chat_history[-6:]:
                role = "User" if msg["role"] == "user" else "Assistant"
                prompt_parts.append(f"{role}: {msg['content']}")
            prompt_parts.append("\n---\n")

        prompt_parts.extend([
            "Context:",
            context,
            f"\nQuestion: {query}",
            "\nProvide a clear, concise answer:",
        ])

        prompt = "\n".join(prompt_parts)

        try:
            response = self.llm.complete(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Error generating answer: {e}"

    # ── Feedback / statistics ─────────────────────────────────────────────────

    def add_feedback(self, query: str, answer: str, is_positive: bool) -> None:
        """
        Add user feedback to the dynamic learning system.

        Args:
            query: Original query.
            answer: Generated answer.
            is_positive: Whether the user found the answer helpful.
        """
        if is_positive:
            print(f"\n{'=' * 70}")
            print(f"👍 FEEDBACK RECEIVED BY LEARNING MANAGER")
            print(f"{'=' * 70}")
            print(f"Query: '{query[:70]}...'")
            print("Action: Successfully routed to LearningManager and saved as Few-Shot Example!")
            
            self.learning_manager.add_good_example(
                query=query,
                answer=answer,
                feedback_score=1.0,
            )
            print(f"{'=' * 70}\n")


